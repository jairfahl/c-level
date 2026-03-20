"""
Knowledge Base Service — Upload, extração de texto e consulta de documentos organizacionais.

Documentos são globais (não vinculados a um caso específico) e categorizados
por domínio financeiro e tipo de decisão.
"""
from __future__ import annotations

import csv
import io
from typing import Optional
from uuid import UUID

from fastapi import UploadFile
from sqlalchemy import case as sa_case, select
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.audit_logger import AuditAction, AuditLogger
from app.core.exceptions import (
    DocumentExtractionError,
    DocumentNotFoundError,
    DocumentTooLargeError,
    UnsupportedFileTypeError,
)
from app.models.enums import DecisionType, FinancialDomain
from app.models.knowledge_document import KnowledgeDocument


class KnowledgeBaseService:
    MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB
    ALLOWED_TYPES = {"pdf", "docx", "txt", "xlsx", "xls", "pptx", "csv", "md"}
    MAX_DOCS_PER_PROMPT = 3
    MAX_CHARS_PER_DOC = 2000
    MAX_XLSX_ROWS = 50_000

    # ── Upload ────────────────────────────────────────────────────────────

    @staticmethod
    async def upload_document(
        session: AsyncSession,
        file: UploadFile,
        title: str,
        financial_domain: FinancialDomain,
        decision_type: Optional[DecisionType] = None,
        description: Optional[str] = None,
        uploaded_by: Optional[str] = None,
    ) -> KnowledgeDocument:
        content = await file.read()
        size = len(content)

        if size > KnowledgeBaseService.MAX_FILE_SIZE:
            raise DocumentTooLargeError(size)

        file_type = KnowledgeBaseService._detect_file_type(file.filename or "")
        extracted_text = KnowledgeBaseService._extract_text(content, file_type)

        doc = KnowledgeDocument(
            title=title,
            description=description,
            original_filename=file.filename or "unknown",
            file_type=file_type,
            file_size_bytes=size,
            extracted_text=extracted_text,
            text_length=len(extracted_text),
            financial_domain=financial_domain,
            decision_type=decision_type,
            uploaded_by=uploaded_by,
        )
        session.add(doc)

        AuditLogger.log(
            session=session,
            action=AuditAction.DOCUMENT_UPLOADED,
            payload={
                "title": title,
                "file_type": file_type,
                "file_size_bytes": size,
                "text_length": len(extracted_text),
                "financial_domain": financial_domain.value,
                "decision_type": decision_type.value if decision_type else None,
            },
        )

        return doc

    # ── File type detection ───────────────────────────────────────────────

    @staticmethod
    def _detect_file_type(filename: str) -> str:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext not in KnowledgeBaseService.ALLOWED_TYPES:
            raise UnsupportedFileTypeError(ext or "(sem extensão)")
        return ext

    # ── Text extraction ───────────────────────────────────────────────────

    @staticmethod
    def _extract_text(content: bytes, file_type: str) -> str:
        extractors = {
            "pdf": KnowledgeBaseService._extract_pdf,
            "docx": KnowledgeBaseService._extract_docx,
            "txt": KnowledgeBaseService._extract_txt,
            "xlsx": KnowledgeBaseService._extract_xlsx,
            "xls": KnowledgeBaseService._extract_xls,
            "pptx": KnowledgeBaseService._extract_pptx,
            "csv": KnowledgeBaseService._extract_csv,
            "md": KnowledgeBaseService._extract_txt,
        }
        extractor = extractors.get(file_type)
        if not extractor:
            raise UnsupportedFileTypeError(file_type)

        text = extractor(content)
        if not text or not text.strip():
            raise DocumentExtractionError("Nenhum texto foi extraído do documento")
        return text.strip()

    @staticmethod
    def _extract_pdf(content: bytes) -> str:
        try:
            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(content))
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages)
        except Exception as e:
            raise DocumentExtractionError(f"PDF: {e}")

    @staticmethod
    def _extract_docx(content: bytes) -> str:
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            raise DocumentExtractionError(f"DOCX: {e}")

    @staticmethod
    def _extract_txt(content: bytes) -> str:
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return content.decode("latin-1")
            except Exception as e:
                raise DocumentExtractionError(f"TXT: {e}")

    @staticmethod
    def _extract_xlsx(content: bytes) -> str:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            lines: list[str] = []
            total_rows = 0
            for sheet in wb.worksheets:
                lines.append(f"[{sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    total_rows += 1
                    if total_rows > KnowledgeBaseService.MAX_XLSX_ROWS:
                        lines.append(f"... (truncado em {KnowledgeBaseService.MAX_XLSX_ROWS} linhas)")
                        wb.close()
                        return "\n".join(lines)
                    cells = [str(c) if c is not None else "" for c in row]
                    lines.append("\t".join(cells))
            wb.close()
            return "\n".join(lines)
        except Exception as e:
            raise DocumentExtractionError(f"XLSX: {e}")

    @staticmethod
    def _extract_xls(content: bytes) -> str:
        try:
            return KnowledgeBaseService._extract_xlsx(content)
        except DocumentExtractionError:
            raise DocumentExtractionError(
                "XLS: Formato .xls (BIFF) não suportado. "
                "Por favor, salve o arquivo como .xlsx e tente novamente."
            )

    @staticmethod
    def _extract_pptx(content: bytes) -> str:
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            slides_text: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                parts: list[str] = [f"[Slide {i}]"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                parts.append(text)
                slides_text.append("\n".join(parts))
            return "\n\n".join(slides_text)
        except Exception as e:
            raise DocumentExtractionError(f"PPTX: {e}")

    @staticmethod
    def _extract_csv(content: bytes) -> str:
        try:
            try:
                text = content.decode("utf-8")
            except UnicodeDecodeError:
                text = content.decode("latin-1")

            try:
                dialect = csv.Sniffer().sniff(text[:4096])
            except csv.Error:
                dialect = None

            reader = csv.reader(io.StringIO(text), dialect=dialect, delimiter=";" if dialect is None else dialect.delimiter)
            lines: list[str] = []
            for row in reader:
                lines.append("\t".join(row))
            return "\n".join(lines)
        except Exception as e:
            raise DocumentExtractionError(f"CSV: {e}")

    # ── Listing ───────────────────────────────────────────────────────────

    @staticmethod
    async def list_documents(
        session: AsyncSession,
        domain: Optional[FinancialDomain] = None,
        decision_type: Optional[DecisionType] = None,
        active_only: bool = True,
    ) -> list[KnowledgeDocument]:
        stmt = select(KnowledgeDocument)
        if active_only:
            stmt = stmt.where(KnowledgeDocument.active == True)
        if domain:
            stmt = stmt.where(KnowledgeDocument.financial_domain == domain)
        if decision_type:
            stmt = stmt.where(KnowledgeDocument.decision_type == decision_type)
        stmt = stmt.order_by(KnowledgeDocument.created_at.desc())
        result = await session.execute(stmt)
        return list(result.scalars().all())

    # ── Get by ID ─────────────────────────────────────────────────────────

    @staticmethod
    async def get_document(
        session: AsyncSession,
        document_id: UUID,
    ) -> KnowledgeDocument:
        result = await session.execute(
            select(KnowledgeDocument).where(KnowledgeDocument.id == document_id)
        )
        doc = result.scalar_one_or_none()
        if doc is None:
            raise DocumentNotFoundError()
        return doc

    # ── Soft-delete ───────────────────────────────────────────────────────

    @staticmethod
    async def delete_document(
        session: AsyncSession,
        document_id: UUID,
    ) -> KnowledgeDocument:
        doc = await KnowledgeBaseService.get_document(session, document_id)
        doc.active = False

        AuditLogger.log(
            session=session,
            action=AuditAction.DOCUMENT_DELETED,
            payload={
                "document_id": str(doc.id),
                "title": doc.title,
            },
        )

        return doc

    # ── Query relevant snippets for LLM injection ────────────────────────

    @staticmethod
    async def query_relevant_snippets(
        session: AsyncSession,
        domain: FinancialDomain,
        decision_type: Optional[DecisionType] = None,
    ) -> list[str]:
        """Retorna até MAX_DOCS_PER_PROMPT snippets relevantes, cada um truncado
        em MAX_CHARS_PER_DOC caracteres. Prioriza documentos com decision_type
        correspondente, depois os genéricos (decision_type IS NULL)."""
        stmt = (
            select(KnowledgeDocument.extracted_text)
            .where(
                KnowledgeDocument.active == True,
                KnowledgeDocument.financial_domain == domain,
            )
        )
        if decision_type:
            stmt = stmt.where(
                (KnowledgeDocument.decision_type == decision_type)
                | (KnowledgeDocument.decision_type.is_(None))
            )
            stmt = stmt.order_by(
                sa_case(
                    (KnowledgeDocument.decision_type == decision_type, 0),
                    else_=1,
                ),
                KnowledgeDocument.created_at.desc(),
            )
        else:
            stmt = stmt.order_by(KnowledgeDocument.created_at.desc())

        stmt = stmt.limit(KnowledgeBaseService.MAX_DOCS_PER_PROMPT)
        result = await session.execute(stmt)
        texts = result.scalars().all()

        return [
            t[: KnowledgeBaseService.MAX_CHARS_PER_DOC]
            for t in texts
        ]
